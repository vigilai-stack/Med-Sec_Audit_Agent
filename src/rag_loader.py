"""
RAG Document Loader
Loads and processes documents for the RAG knowledge base

Supports:
- Local files: PDF, TXT, MD, DOCX, PPTX, CSV, XLSX
- Web URLs:    listed in knowledge_base/urls.txt (one URL per line)
- API sources: listed in knowledge_base/api_sources.json (placeholder stubs)
"""

import os
import json
from typing import List, Dict, Any, Optional

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, TextLoader, WebBaseLoader

from src.config import config


class RAGDocumentLoader:
    """
    Loads documents from the knowledge_base directory
    
    Supports:
    - PDF files (.pdf)
    - Text files (.txt)
    - Markdown files (.md)
    - Word documents (.docx)        — requires python-docx
    - PowerPoint files (.pptx)      — requires python-pptx
    - CSV files (.csv)              — requires pandas
    - Excel files (.xlsx)           — requires openpyxl
    - Web URLs                      — listed in knowledge_base/urls.txt
    - API sources                   — listed in knowledge_base/api_sources.json
    """
    
    def __init__(self, knowledge_base_path: Optional[str] = None):
        self.knowledge_base_path = knowledge_base_path or os.path.join(
            config.base_path, "knowledge_base", "documents"
        )
        # urls.txt and api_sources.json live one level up (in knowledge_base/)
        self.knowledge_base_root = os.path.dirname(self.knowledge_base_path)
        self.urls_file = os.path.join(self.knowledge_base_root, "urls.txt")
        self.api_sources_file = os.path.join(self.knowledge_base_root, "api_sources.json")
        self.supported_extensions = ['.pdf', '.txt', '.md', '.docx', '.pptx', '.csv', '.xlsx']
        self.documents = []
        
        print(f"📚 RAG Document Loader initialized")
        print(f"   Knowledge base: {self.knowledge_base_path}")
        print(f"   URLs file: {self.urls_file}")
        print(f"   API sources: {self.api_sources_file}")
        print(f"   Supported extensions: {', '.join(self.supported_extensions)}")
    
    def load_all_documents(self) -> List[Document]:
        """
        Load all documents from the knowledge base (files + URLs + API sources)
        
        Returns:
            List of Document objects
        """
        all_docs = []
        
        # ── 1. Load local files ──
        if not os.path.exists(self.knowledge_base_path):
            print(f"⚠️ Knowledge base path not found: {self.knowledge_base_path}")
        else:
            for root, dirs, files in os.walk(self.knowledge_base_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    ext = os.path.splitext(file)[1].lower()
                    
                    if ext == '.pdf':
                        docs = self._load_pdf(file_path)
                    elif ext in ['.txt', '.md']:
                        docs = self._load_text(file_path)
                    elif ext == '.docx':
                        docs = self._load_docx(file_path)
                    elif ext == '.pptx':
                        docs = self._load_pptx(file_path)
                    elif ext == '.csv':
                        docs = self._load_csv(file_path)
                    elif ext == '.xlsx':
                        docs = self._load_xlsx(file_path)
                    else:
                        print(f"   ⏭️  Skipping unsupported file: {file}")
                        continue
                    
                    all_docs.extend(docs)
                    print(f"   ✅ Loaded: {file} ({len(docs)} chunks)")
        
        # ── 2. Load web URLs ──
        web_docs = self._load_urls()
        if web_docs:
            all_docs.extend(web_docs)
        
        # ── 3. Load API sources ──
        api_docs = self._load_api_sources()
        if api_docs:
            all_docs.extend(api_docs)
        
        self.documents = all_docs
        print(f"\n✅ Total documents loaded: {len(all_docs)}")
        return all_docs
    
    def _load_pdf(self, file_path: str) -> List[Document]:
        """Load a PDF file"""
        try:
            loader = PyPDFLoader(file_path)
            return loader.load()
        except Exception as e:
            print(f"   ⚠️ Error loading PDF {file_path}: {e}")
            return []
    
    def _load_text(self, file_path: str) -> List[Document]:
        """Load a text or markdown file"""
        try:
            loader = TextLoader(file_path, encoding='utf-8')
            return loader.load()
        except Exception as e:
            print(f"   ⚠️ Error loading text {file_path}: {e}")
            return []
    
    def _load_docx(self, file_path: str) -> List[Document]:
        """Load a Word .docx file using python-docx directly"""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            print("   ⚠️ python-docx not installed — skipping .docx files. Install with: pip install python-docx")
            return []
        try:
            doc = DocxDocument(file_path)
            # Combine all paragraphs into a single text block
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        paragraphs.append(" | ".join(cells))
            content = "\n".join(paragraphs)
            if not content.strip():
                return []
            return [Document(
                page_content=content,
                metadata={"source": file_path, "type": "docx"}
            )]
        except Exception as e:
            print(f"   ⚠️ Error loading DOCX {file_path}: {e}")
            return []
    
    def _load_pptx(self, file_path: str) -> List[Document]:
        """Load a PowerPoint .pptx file using python-pptx directly"""
        try:
            from pptx import Presentation
        except ImportError:
            print("   ⚠️ python-pptx not installed — skipping .pptx files. Install with: pip install python-pptx")
            return []
        try:
            prs = Presentation(file_path)
            docs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from all shapes in the slide
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    content = "\n".join(texts)
                    docs.append(Document(
                        page_content=content,
                        metadata={"source": file_path, "slide": slide_num, "type": "pptx"}
                    ))
            return docs
        except Exception as e:
            print(f"   ⚠️ Error loading PPTX {file_path}: {e}")
            return []
    
    def _load_csv(self, file_path: str) -> List[Document]:
        """Load a CSV file (one Document per row)"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            docs = []
            for idx, row in df.iterrows():
                content = "\n".join(f"{col}: {val}" for col, val in row.items())
                docs.append(Document(
                    page_content=content,
                    metadata={"source": file_path, "row": idx}
                ))
            return docs
        except Exception as e:
            print(f"   ⚠️ Error loading CSV {file_path}: {e}")
            return []
    
    def _load_xlsx(self, file_path: str) -> List[Document]:
        """Load an Excel .xlsx file (one Document per sheet)"""
        try:
            import pandas as pd
            xlsx = pd.ExcelFile(file_path)
            docs = []
            for sheet_name in xlsx.sheet_names:
                df = xlsx.parse(sheet_name)
                content = df.to_string(index=False)
                docs.append(Document(
                    page_content=f"Sheet: {sheet_name}\n\n{content}",
                    metadata={"source": file_path, "sheet": sheet_name}
                ))
            return docs
        except Exception as e:
            print(f"   ⚠️ Error loading XLSX {file_path}: {e}")
            return []
    
    # ═══════════════════════════════════════════════════════════════════════
    # WEB URL LOADING
    # ═══════════════════════════════════════════════════════════════════════
    
    def _load_urls(self) -> List[Document]:
        """
        Load web pages from URLs listed in knowledge_base/urls.txt
        
        Format of urls.txt (one URL per line, # for comments):
            # HIPAA resources
            https://www.hhs.gov/hipaa/index.html
            # NIST
            https://www.nist.gov/cyberframework
        """
        if not os.path.exists(self.urls_file):
            return []
        
        # Read URLs (skip comments and blank lines)
        with open(self.urls_file, 'r', encoding='utf-8') as f:
            urls = []
            for line in f:
                line = line.strip()
                if line and not line.startswith('#'):
                    urls.append(line)
        
        if not urls:
            return []
        
        print(f"\n🌐 Loading {len(urls)} web URLs...")
        all_docs = []
        
        for url in urls:
            docs = self._load_web_page(url)
            all_docs.extend(docs)
            print(f"   ✅ Loaded: {url} ({len(docs)} chunks)")
        
        return all_docs
    
    def _load_web_page(self, url: str) -> List[Document]:
        """Load a single web page"""
        try:
            loader = WebBaseLoader(url)
            docs = loader.load()
            # Add source metadata
            for doc in docs:
                doc.metadata["source"] = url
                doc.metadata["type"] = "web"
            return docs
        except Exception as e:
            print(f"   ⚠️ Error loading URL {url}: {e}")
            return []
    
    # ═══════════════════════════════════════════════════════════════════════
    # API SOURCE LOADING (placeholders — fill in with real API calls)
    # ═══════════════════════════════════════════════════════════════════════
    
    def _load_api_sources(self) -> List[Document]:
        """
        Load data from API sources defined in knowledge_base/api_sources.json
        
        Format of api_sources.json:
        [
            {
                "name": "PubMed",
                "type": "pubmed",
                "query": "HIPAA compliance healthcare",
                "max_results": 10,
                "api_key": ""           ← optional, leave empty for no auth
            },
            {
                "name": "openFDA",
                "type": "openfda",
                "endpoint": "recalls",
                "search": "device",
                "limit": 20
            },
            {
                "name": "ClinicalTrials.gov",
                "type": "clinical_trials",
                "condition": "cardiovascular",
                "max_results": 10
            }
        ]
        """
        if not os.path.exists(self.api_sources_file):
            return []
        
        try:
            with open(self.api_sources_file, 'r', encoding='utf-8') as f:
                sources = json.load(f)
        except Exception as e:
            print(f"   ⚠️ Error reading api_sources.json: {e}")
            return []
        
        if not sources:
            return []
        
        print(f"\n🔌 Loading {len(sources)} API sources...")
        all_docs = []
        
        for source in sources:
            source_type = source.get("type", "")
            source_name = source.get("name", source_type)
            
            if source_type == "pubmed":
                docs = self._fetch_pubmed(source)
            elif source_type == "openfda":
                docs = self._fetch_openfda(source)
            elif source_type == "clinical_trials":
                docs = self._fetch_clinical_trials(source)
            elif source_type == "rxnorm":
                docs = self._fetch_rxnorm(source)
            elif source_type == "dailymed":
                docs = self._fetch_dailymed(source)
            elif source_type == "nih_reporter":
                docs = self._fetch_nih_reporter(source)
            elif source_type == "medlineplus":
                docs = self._fetch_medlineplus(source)
            else:
                print(f"   ⏭️  Unknown API type '{source_type}' for '{source_name}'")
                continue
            
            all_docs.extend(docs)
            print(f"   ✅ Loaded from {source_name}: {len(docs)} documents")
        
        return all_docs
    
    def _fetch_pubmed(self, source: Dict) -> List[Document]:
        """
        Fetch articles from PubMed (NCBI E-utilities API)
        
        Setup:
            1. Get an NCBI API key (optional but recommended): https://www.ncbi.nlm.nih.gov/account/
            2. Set it in api_sources.json as "api_key": "your_key_here"
            3. Or leave empty for slower rate-limited access
        
        source config:
            {
                "name": "PubMed",
                "type": "pubmed",
                "query": "HIPAA compliance healthcare",
                "max_results": 10,
                "api_key": ""
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for PubMed API")
            return []
        
        try:
            query = source.get("query", "")
            max_results = source.get("max_results", 10)
            api_key = source.get("api_key", "")
            
            if not query:
                print("   ⚠️ PubMed source missing 'query'")
                return []
            
            # Step 1: Search PubMed for article IDs
            search_url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi"
            params = {
                "db": "pubmed",
                "term": query,
                "retmax": max_results,
                "retmode": "json"
            }
            if api_key:
                params["api_key"] = api_key
            
            resp = requests.get(search_url, params=params, timeout=30)
            resp.raise_for_status()
            pmids = resp.json().get("esearchresult", {}).get("idlist", [])
            
            if not pmids:
                print(f"   ℹ️  No PubMed results for query: {query}")
                return []
            
            # Step 2: Fetch summaries for those IDs
            fetch_url = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esummary.fcgi"
            params = {
                "db": "pubmed",
                "id": ",".join(pmids),
                "retmode": "json"
            }
            if api_key:
                params["api_key"] = api_key
            
            resp = requests.get(fetch_url, params=params, timeout=30)
            resp.raise_for_status()
            result = resp.json().get("result", {})
            
            docs = []
            for pmid in pmids:
                article = result.get(pmid, {})
                title = article.get("title", "")
                authors = ", ".join(
                    a.get("name", "") for a in article.get("authors", [])
                )
                journal = article.get("fulljournalname", "")
                pubdate = article.get("pubdate", "")
                
                content = f"Title: {title}\nAuthors: {authors}\nJournal: {journal}\nPublished: {pubdate}"
                
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"PubMed:{pmid}",
                        "type": "pubmed",
                        "pmid": pmid,
                        "title": title,
                        "url": f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/"
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching PubMed: {e}")
            return []
    
    def _fetch_openfda(self, source: Dict) -> List[Document]:
        """
        Fetch data from openFDA API (no API key required, but rate-limited)
        
        Setup:
            1. No setup needed for basic access (240 requests/minute per IP)
            2. For higher limits, get a key: https://open.fda.gov/api/reference/
            3. Set it in api_sources.json as "api_key": "your_key_here"
        
        source config:
            {
                "name": "openFDA Recalls",
                "type": "openfda",
                "endpoint": "recalls",          ← recalls, device, drug, food, etc.
                "search": "device",             ← search term
                "limit": 20
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for openFDA API")
            return []
        
        try:
            endpoint = source.get("endpoint", "recalls")
            search = source.get("search", "")
            limit = source.get("limit", 20)
            api_key = source.get("api_key", "")
            
            url = f"https://api.fda.gov/{endpoint}.json"
            params = {"limit": limit}
            if search:
                params["search"] = search
            if api_key:
                params["api_key"] = api_key
            
            resp = requests.get(url, params=params, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            
            results = data.get("results", [])
            docs = []
            for item in results:
                # Convert JSON to readable text
                content = json.dumps(item, indent=2)
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"openFDA:{endpoint}",
                        "type": "openfda",
                        "endpoint": endpoint
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching openFDA: {e}")
            return []
    
    def _fetch_clinical_trials(self, source: Dict) -> List[Document]:
        """
        Fetch clinical trials from ClinicalTrials.gov API v2
        No API key required.
        
        source config:
            {
                "name": "ClinicalTrials.gov",
                "type": "clinical_trials",
                "condition": "cardiovascular",
                "max_results": 10
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for ClinicalTrials.gov API")
            return []
        
        try:
            condition = source.get("condition", "")
            max_results = source.get("max_results", 10)
            
            url = "https://clinicaltrials.gov/api/v2/studies"
            params = {
                "query.term": condition,
                "pageSize": max_results,
                "format": "json"
            }
            
            resp = requests.get(url, params=params, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            
            studies = data.get("studies", [])
            docs = []
            for study in studies:
                protocol = study.get("protocolSection", {})
                id_module = protocol.get("identificationModule", {})
                
                nct_id = id_module.get("nctId", "")
                title = id_module.get("officialTitle", id_module.get("briefTitle", ""))
                status = protocol.get("statusModule", {}).get("overallStatus", "")
                
                content = f"NCT ID: {nct_id}\nTitle: {title}\nStatus: {status}"
                
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"ClinicalTrials:{nct_id}",
                        "type": "clinical_trials",
                        "nct_id": nct_id,
                        "url": f"https://clinicaltrials.gov/study/{nct_id}"
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching ClinicalTrials.gov: {e}")
            return []
    
    def _fetch_rxnorm(self, source: Dict) -> List[Document]:
        """
        Fetch drug information from RxNorm API (NLM).
        No API key required. Rate limit: 20 requests/second per IP.
        
        Uses approximateTerm search (handles both drug names and general terms).
        Then fetches full drug info for each RxCUI found.
        
        source config:
            {
                "name": "RxNorm - Aspirin",
                "type": "rxnorm",
                "query": "aspirin",
                "max_results": 25
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for RxNorm API")
            return []
        
        try:
            query = source.get("query", "")
            max_results = source.get("max_results", 25)
            
            if not query:
                print("   ⚠️ RxNorm source missing 'query'")
                return []
            
            # Step 1: Approximate term search to get RxCUIs
            search_url = "https://rxnav.nlm.nih.gov/REST/approximateTerm.json"
            params = {"term": query, "maxEntries": max_results}
            
            resp = requests.get(search_url, params=params, timeout=30)
            resp.raise_for_status()
            candidates = resp.json().get("approximateGroup", {}).get("candidate", [])
            
            if not candidates:
                print(f"   ℹ️  No RxNorm results for: {query}")
                return []
            
            # Deduplicate by rxcui
            seen_rxcuis = set()
            unique_candidates = []
            for c in candidates:
                rxcui = c.get("rxcui", "")
                if rxcui and rxcui not in seen_rxcuis:
                    seen_rxcuis.add(rxcui)
                    unique_candidates.append(c)
            
            # Step 2: Get drug properties for each unique RxCUI
            docs = []
            for cand in unique_candidates[:max_results]:
                rxcui = cand.get("rxcui", "")
                name = cand.get("name", "")
                source_name = cand.get("source", "")
                score = cand.get("score", "")
                
                content = f"Drug: {name}\nRxCUI: {rxcui}\nSource: {source_name}\nMatch Score: {score}"
                
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"RxNorm:{rxcui}",
                        "type": "rxnorm",
                        "rxcui": rxcui,
                        "drug_name": name,
                        "url": f"https://mor.nlm.nih.gov/RxNav/search?searchBy=RXCUI&searchArg={rxcui}"
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching RxNorm: {e}")
            return []
    
    def _fetch_dailymed(self, source: Dict) -> List[Document]:
        """
        Fetch official drug labeling from DailyMed (FDA/NLM).
        No API key required.
        
        source config:
            {
                "name": "DailyMed - Drug Labels",
                "type": "dailymed",
                "query": "aspirin",
                "max_results": 20
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for DailyMed API")
            return []
        
        try:
            query = source.get("query", "")
            max_results = source.get("max_results", 20)
            
            if not query:
                print("   ⚠️ DailyMed source missing 'query'")
                return []
            
            # Search for SPL documents
            url = "https://dailymed.nlm.nih.gov/dailymed/services/v2/spls.json"
            params = {"drug_name": query, "pagesize": max_results}
            
            resp = requests.get(url, params=params, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            
            spls = data.get("data", [])
            docs = []
            for spl in spls:
                setid = spl.get("setid", "")
                title = spl.get("title", "")
                published = spl.get("published_date", "")
                
                content = f"Drug Label: {title}\nSet ID: {setid}\nPublished: {published}"
                
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"DailyMed:{setid}",
                        "type": "dailymed",
                        "setid": setid,
                        "url": f"https://dailymed.nlm.nih.gov/dailymed/drugInfo.cfm?setid={setid}"
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching DailyMed: {e}")
            return []
    
    def _fetch_nih_reporter(self, source: Dict) -> List[Document]:
        """
        Fetch funded research projects from NIH RePORTER API.
        No API key required.
        
        source config:
            {
                "name": "NIH RePORTER - Healthcare Cybersecurity",
                "type": "nih_reporter",
                "query": "healthcare cybersecurity HIPAA",
                "max_results": 20
            }
        """
        try:
            import requests
        except ImportError:
            print("   ⚠️ requests library needed for NIH RePORTER API")
            return []
        
        try:
            query = source.get("query", "")
            max_results = source.get("max_results", 20)
            
            if not query:
                print("   ⚠️ NIH RePORTER source missing 'query'")
                return []
            
            url = "https://api.reporter.nih.gov/v2/projects/search"
            payload = {
                "criteria": {
                    "text": query
                },
                "limit": max_results,
                "offset": 0
            }
            
            resp = requests.post(url, json=payload, timeout=30)
            resp.raise_for_status()
            data = resp.json()
            
            results = data.get("results", [])
            docs = []
            for project in results:
                project_id = project.get("project_num", "")
                title = project.get("project_title", "")
                abstract = project.get("abstract_text", "")
                agency_info = project.get("agency_ic_funding")
                agency = ""
                if agency_info and isinstance(agency_info, dict):
                    agency = agency_info.get("abbreviation", "")
                fiscal_year = project.get("fiscal_year", "")
                pi_list = project.get("principal_investigators", [])
                pi_name = pi_list[0].get("full_name", "") if pi_list else ""
                
                content = f"Project: {title}\nProject ID: {project_id}\nAgency: {agency}\nFY: {fiscal_year}\nPI: {pi_name}\nAbstract: {abstract[:500]}"
                
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"NIHRePORTER:{project_id}",
                        "type": "nih_reporter",
                        "project_id": project_id,
                        "title": title
                    }
                ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching NIH RePORTER: {e}")
            return []
    
    def _fetch_medlineplus(self, source: Dict) -> List[Document]:
        """
        Fetch health topic information from MedlinePlus (NLM).
        No API key required. Returns XML, parsed internally.
        
        source config:
            {
                "name": "MedlinePlus - HIPAA",
                "type": "medlineplus",
                "query": "HIPAA",
                "max_results": 20
            }
        """
        try:
            import requests
            import xml.etree.ElementTree as ET
            import re
        except ImportError:
            print("   ⚠️ requests library needed for MedlinePlus API")
            return []
        
        try:
            query = source.get("query", "")
            max_results = source.get("max_results", 20)
            
            if not query:
                print("   ⚠️ MedlinePlus source missing 'query'")
                return []
            
            # Search MedlinePlus health topics (returns XML)
            url = "https://wsearch.nlm.nih.gov/wssearch"
            params = {
                "db": "healthTopics",
                "term": query,
                "retmax": max_results
            }
            
            resp = requests.get(url, params=params, timeout=30)
            resp.raise_for_status()
            
            root = ET.fromstring(resp.text)
            docs = []
            
            for doc_el in root.findall(".//document"):
                doc_url = doc_el.get("url", "")
                title = ""
                summary = ""
                
                # Content elements have name attributes (title, FullSummary, etc.)
                for content in doc_el.findall("content"):
                    name = content.get("name", "")
                    if name == "title" and content.text:
                        title = content.text.strip()
                    elif name == "FullSummary" and content.text:
                        # Strip HTML tags
                        summary = re.sub(r"<[^>]+>", "", content.text).strip()
                
                if title or summary:
                    content_text = f"Health Topic: {title}\nSummary: {summary}"
                    docs.append(Document(
                        page_content=content_text,
                        metadata={
                            "source": f"MedlinePlus:{title}",
                            "type": "medlineplus",
                            "url": doc_url
                        }
                    ))
            
            return docs
        except Exception as e:
            print(f"   ⚠️ Error fetching MedlinePlus: {e}")
            return []
    
    # ═══════════════════════════════════════════════════════════════════════
    
    def load_from_texts(self, texts: List[str], metadata: Optional[List[Dict]] = None) -> List[Document]:
        """Load documents from text strings"""
        documents = []
        for i, text in enumerate(texts):
            meta = metadata[i] if metadata else {"source": f"text_{i}"}
            documents.append(Document(page_content=text, metadata=meta))
        return documents
    
    def get_document_summary(self) -> Dict[str, Any]:
        """Get summary of loaded documents"""
        return {
            "total_documents": len(self.documents),
            "knowledge_base_path": self.knowledge_base_path
        }