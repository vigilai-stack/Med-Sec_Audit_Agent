# ==============================================================================
# Report Generation Agent — Professional Multi-Format Audit Reports
# ==============================================================================
# Generates enterprise-grade reports in PDF, Word, Excel, PowerPoint, HTML
# Uses reportlab, python-docx, openpyxl, python-pptx
# ==============================================================================

import os
import json
import base64
import tempfile
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns

try:
    from IPython.display import display, HTML, Markdown
except ImportError:
    def display(value=None, *args, **kwargs):
        return value
    class HTML(str):
        pass
    class Markdown(str):
        pass

# ── DOCX ──────────────────────────────────────────────────────────────────────
DOCX_AVAILABLE = True
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
except ImportError:
    DOCX_AVAILABLE = False


class ReportGenerationAgent:
    """
    Multi-agent report generator producing enterprise-grade audit reports.

    Pipeline:
    1. Critique Agent  — analyzes results, extracts key insights
    2. Visualization Agent — creates charts & diagrams (matplotlib)
    3. Writer Agent — drafts structured narrative
    4. Render Agent — exports to PDF, DOCX, XLSX, PPTX, HTML
    """

    # ── Color Palette (shared across all formats) ──────────────────────────
    NAVY   = "#12355B"
    TEAL   = "#0B7A75"
    RED    = "#C93636"
    GOLD   = "#B7791F"
    GREEN  = "#13795B"
    SOFT   = "#F5F7FA"
    LINE   = "#D8DEE9"
    INK    = "#172033"
    MUTED  = "#667085"
    WHITE  = "#FFFFFF"
    LIGHT_GOLD = "#F5E6B8"
    LIGHT_TEAL = "#BDE7E3"

    def __init__(self, workspace_dir: str = "./medsec_sandbox", skills_path: str = None):
        self.workspace = workspace_dir
        self.skills_path = skills_path or "./.agents/skills/"
        self.report_data = {}
        self.audit_results = None

        self.reports_dir = os.path.join(workspace_dir, "reports")
        os.makedirs(self.reports_dir, exist_ok=True)

        print("📊 ReportGenerationAgent initialized")
        print(f"   Workspace: {self.workspace}")
        print(f"   Reports Directory: {self.reports_dir}")

    # ════════════════════════════════════════════════════════════════════════
    # PUBLIC API
    # ════════════════════════════════════════════════════════════════════════

    def generate_report(self, audit_results: Dict[str, Any],
                        report_format: str = "all") -> Dict[str, str]:
        """Generate complete report from audit results."""
        self.audit_results = audit_results
        self.report_data = self._prepare_report_data()

        print("\n" + "=" * 70)
        print("📊 REPORT GENERATION AGENT — Starting Pipeline")
        print("=" * 70)

        print("\n  Phase 1: Critique Agent — Analyzing audit results...")
        critique = self._run_critique_agent()

        print("\n  Phase 2: Visualization Agent — Creating charts...")
        visuals = self._run_visualization_agent()

        print("\n  Phase 3: Writer Agent — Drafting report...")
        narrative = self._build_narrative(critique)

        print("\n  Phase 4: Render Agent — Exporting to all formats...")
        outputs = self._run_render_agent(narrative, visuals, critique)

        print("\n✅ Report generation complete!")
        for fmt, path in outputs.items():
            print(f"   {fmt.upper()}: {path}")
        return outputs

    def display_report(self, report_paths: Dict[str, str]):
        """Display generated report in Colab/Notebook."""
        print("\n" + "=" * 70)
        print("📊 REPORT GENERATED SUCCESSFULLY")
        print("=" * 70)
        for fmt, path in report_paths.items():
            if os.path.exists(path):
                print(f"  {fmt.upper()}: {path}")
        html_path = report_paths.get("html")
        if html_path and os.path.exists(html_path):
            with open(html_path, "r", encoding="utf-8") as f:
                display(HTML(f.read()))

    # ════════════════════════════════════════════════════════════════════════
    # PHASE 1: PREPARE DATA
    # ════════════════════════════════════════════════════════════════════════

    def _prepare_report_data(self) -> Dict[str, Any]:
        """Extract and structure data from audit results."""
        data = {
            "audit_id": self.audit_results.get("audit_id", "UNKNOWN"),
            "timestamp": self.audit_results.get("timestamp", datetime.now().isoformat()),
            "duration": self.audit_results.get("duration_seconds", 0),
            "risk_assessment": self.audit_results.get("risk_assessment", {}),
            "remediation": self.audit_results.get("remediation", {}),
            "compliance": self.audit_results.get("compliance", {}),
            "recommendations": self.audit_results.get("recommendations", []),
            "report_context": self.audit_results.get("report_context", {}),
            "red_team": {}, "blue_team": {}, "green_team": {},
            "compliance_details": {}, "baseline": {},
            "phi_anonymization": {},
            "data_processed": self.audit_results.get("data_processed", {}),
        }
        raw = self.audit_results.get("raw_results", {})
        data["baseline"] = raw.get("baseline", {})
        data["phi_anonymization"] = raw.get("phi_anonymization", {})

        red = raw.get("red_team", {})
        data["red_team"] = {
            "target": red.get("target", "Unknown"),
            "attack_types": red.get("attack_types", []),
            "overall_risk": red.get("overall_risk", {}),
            "results": red.get("results", {}),
        }
        blue = raw.get("blue_team", {})
        data["blue_team"] = {
            "threat_count": blue.get("threat_count", 0),
            "risk_level": blue.get("risk_level", "LOW"),
            "anomalies": blue.get("anomalies_found", []),
            "runtime_ms": blue.get("runtime_ms", 0),
        }
        green = raw.get("green_team", [])
        data["green_team"] = {"fixes_applied": len(green), "fixes": green}
        compliance = raw.get("compliance", {})
        data["compliance_details"] = {
            "standard": compliance.get("standard", "HIPAA"),
            "score": compliance.get("compliance_score", 0),
            "passed_checks": compliance.get("passed_checks", 0),
            "total_checks": compliance.get("total_checks", 0),
            "issues": compliance.get("issues_found", []),
            "compliant": compliance.get("compliant", False),
        }
        return data

    # ════════════════════════════════════════════════════════════════════════
    # PHASE 2: CRITIQUE AGENT
    # ════════════════════════════════════════════════════════════════════════

    def _run_critique_agent(self) -> Dict[str, Any]:
        data = self.report_data
        risk_level = data["risk_assessment"].get("overall_risk", "UNKNOWN")
        threat_count = data["risk_assessment"].get("threats_detected", 0)
        compliance_score = data["compliance"].get("score", 0)

        critique = {
            "summary": (
                f"The Med-Sec Audit Agent completed a comprehensive security audit "
                f"of the Healthcare Clinical Documentation System. The audit identified "
                f"a {risk_level} risk level with {threat_count} threats detected. "
                f"HIPAA compliance score: {compliance_score:.1f}%. "
                f"{'The system is HIPAA compliant.' if compliance_score >= 80 else 'The system requires remediation to meet HIPAA standards.'}"
            ),
            "highlights": [],
            "critical_findings": [],
            "improvement_areas": [],
        }

        if data["green_team"]["fixes_applied"] > 0:
            critique["highlights"].append(f"{data['green_team']['fixes_applied']} vulnerabilities automatically remediated")
        if compliance_score >= 80:
            critique["highlights"].append("HIPAA compliance requirements met")
        if data["blue_team"]["threat_count"] > 0:
            critique["highlights"].append(f"{data['blue_team']['threat_count']} threats detected and logged")

        for anomaly in data["blue_team"]["anomalies"][:3]:
            critique["critical_findings"].append(str(anomaly))
        for issue in data["compliance_details"]["issues"][:2]:
            critique["critical_findings"].append(str(issue))
        for rec in data["recommendations"]:
            critique["improvement_areas"].append(str(rec))

        critique_path = os.path.join(self.reports_dir, "critique.json")
        with open(critique_path, "w", encoding="utf-8") as f:
            json.dump(critique, f, indent=2, default=str)
        return critique

    # ════════════════════════════════════════════════════════════════════════
    # PHASE 3: VISUALIZATION AGENT
    # ════════════════════════════════════════════════════════════════════════

    def _run_visualization_agent(self) -> Dict[str, str]:
        visuals = {}
        for name, fn in [
            ("compliance_dashboard", self._create_compliance_dashboard),
            ("risk_chart", self._create_risk_chart),
            ("attack_chart", self._create_attack_chart),
            ("team_summary", self._create_team_summary_chart),
            ("threat_heatmap", self._create_threat_heatmap),
        ]:
            path = fn()
            if path:
                visuals[name] = path
        return visuals

    def _create_compliance_dashboard(self) -> Optional[str]:
        try:
            fig, axes = plt.subplots(1, 2, figsize=(12, 5))
            score = self.report_data["compliance"]["score"]
            axes[0].barh(["HIPAA"], [score / 100], color="#44aaff" if score >= 80 else "#ffaa44")
            axes[0].set_xlim(0, 1)
            axes[0].axvline(x=0.8, color="red", linestyle="--", label="HIPAA Threshold (80%)")
            axes[0].set_title(f"HIPAA Compliance Score: {score:.1f}%")
            axes[0].legend()

            passed = self.report_data["compliance_details"]["passed_checks"]
            total = self.report_data["compliance_details"]["total_checks"]
            axes[1].bar(["Passed", "Failed"], [passed, max(total - passed, 0)], color=["#4CAF50", "#f44336"])
            axes[1].set_title("Compliance Checks")
            axes[1].set_ylabel("Count")
            for i, v in enumerate([passed, max(total - passed, 0)]):
                axes[1].text(i, v + 0.1, str(v), ha="center", va="bottom")

            plt.tight_layout()
            path = os.path.join(self.reports_dir, "compliance_dashboard.png")
            plt.savefig(path, dpi=150, bbox_inches="tight")
            plt.close()
            return path
        except Exception as e:
            print(f"   Compliance dashboard failed: {e}")
            return None

    def _create_risk_chart(self) -> Optional[str]:
        try:
            fig, ax = plt.subplots(figsize=(6, 4))
            risk_levels = ["LOW", "MEDIUM", "HIGH", "CRITICAL"]
            risk_map = {"LOW": 1, "MEDIUM": 2, "HIGH": 3, "CRITICAL": 4}
            current_risk = self.report_data["risk_assessment"].get("overall_risk", "LOW")
            risk_val = risk_map.get(current_risk, 1)
            colors = ["#4CAF50", "#FFC107", "#FF9800", "#f44336"]
            bars = ax.bar(risk_levels, [1, 2, 3, 4], color=["#e8f5e9", "#fff3e0", "#ffe0b2", "#ffcdd2"])
            bars[risk_val - 1].set_color(colors[risk_val - 1])
            bars[risk_val - 1].set_edgecolor("black")
            bars[risk_val - 1].set_linewidth(2)
            ax.set_title(f"Risk Level: {current_risk}")
            ax.set_ylabel("Severity")
            ax.set_ylim(0, 5)
            plt.tight_layout()
            path = os.path.join(self.reports_dir, "risk_chart.png")
            plt.savefig(path, dpi=150, bbox_inches="tight")
            plt.close()
            return path
        except Exception as e:
            print(f"   Risk chart failed: {e}")
            return None

    def _create_attack_chart(self) -> Optional[str]:
        try:
            red_results = self.report_data.get("red_team", {}).get("results", {})
            if not red_results:
                return None
            fig, ax = plt.subplots(figsize=(8, 4))
            attack_names = list(red_results.keys())
            attack_rates = [red_results[a].get("success_rate", 0) * 100 for a in attack_names]
            colors = ["#f44336" if r > 25 else "#FFC107" if r > 10 else "#4CAF50" for r in attack_rates]
            bars = ax.bar(attack_names, attack_rates, color=colors)
            ax.set_ylabel("Success Rate (%)")
            ax.set_title("Attack Success Rates")
            ax.axhline(y=25, color="red", linestyle="--", label="High Risk Threshold (25%)")
            ax.legend()
            for bar, rate in zip(bars, attack_rates):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1, f"{rate:.0f}%", ha="center", va="bottom")
            plt.tight_layout()
            path = os.path.join(self.reports_dir, "attack_chart.png")
            plt.savefig(path, dpi=150, bbox_inches="tight")
            plt.close()
            return path
        except Exception as e:
            print(f"   Attack chart failed: {e}")
            return None

    def _create_team_summary_chart(self) -> Optional[str]:
        try:
            data = self.report_data
            fig, ax = plt.subplots(figsize=(8, 4))
            categories = ["Red Team\nAttacks", "Blue Team\nThreats", "Green Team\nFixes", "Compliance\nScore"]
            red_count = len(data["red_team"].get("results", {}))
            blue_count = data["blue_team"].get("threat_count", 0)
            green_count = data["green_team"].get("fixes_applied", 0)
            comp_score = data["compliance"].get("score", 0)
            values = [red_count, blue_count, green_count, comp_score]
            colors = ["#EF5350", "#42A5F5", "#66BB6A", "#FFA726"]
            bars = ax.bar(categories, values, color=colors)
            ax.set_title("Security Team Summary")
            ax.set_ylabel("Count")
            for bar, val in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5, str(int(val)), ha="center", va="bottom")
            plt.tight_layout()
            path = os.path.join(self.reports_dir, "team_summary.png")
            plt.savefig(path, dpi=150, bbox_inches="tight")
            plt.close()
            return path
        except Exception as e:
            print(f"   Team summary chart failed: {e}")
            return None

    def _create_threat_heatmap(self) -> Optional[str]:
        """Create threat analysis heatmap — large size for maximum readability"""
        try:
            red_results = self.report_data.get("red_team", {}).get("results", {})
            if not red_results:
                return None
            
            # Large size for excellent readability
            fig, ax = plt.subplots(figsize=(11, 5.5))
            attack_names = list(red_results.keys())
            attempts = [red_results[a].get("total_attempts", 0) for a in attack_names]
            successful = [red_results[a].get("successful_attacks", 0) for a in attack_names]
            failed = [a - s for a, s in zip(attempts, successful)]
            
            df = pd.DataFrame({"Attack Type": attack_names, "Successful": successful, "Failed": failed})
            df = df.set_index("Attack Type")
            
            # Create heatmap with better formatting
            sns.heatmap(df.T, annot=True, fmt="d", cmap="RdYlGn_r", ax=ax, 
                       cbar_kws={"label": "Count", "shrink": 0.7},
                       annot_kws={"size": 12, "weight": "bold"})
            ax.set_title("Threat Analysis Heatmap", fontsize=14, fontweight="bold", pad=15)
            
            # Rotate x-axis labels for better readability
            ax.set_xticklabels(ax.get_xticklabels(), rotation=15, ha="right", fontsize=11)
            ax.set_yticklabels(ax.get_yticklabels(), fontsize=11)
            
            plt.tight_layout()
            
            path = os.path.join(self.reports_dir, "threat_heatmap.png")
            plt.savefig(path, dpi=150, bbox_inches="tight")
            plt.close()
            return path
        except Exception as e:
            print(f"   Threat heatmap failed: {e}")
            return None

    # ════════════════════════════════════════════════════════════════════════
    # PHASE 4: WRITER AGENT (narrative)
    # ════════════════════════════════════════════════════════════════════════

    def _build_narrative(self, critique: Dict) -> str:
        data = self.report_data
        risk = data.get("risk_assessment", {})
        compliance = data.get("compliance", {})
        remediation = data.get("remediation", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})

        lines = [
            "# Med-Sec Audit Agent — Security Audit Report",
            "",
            f"**Audit ID:** {data.get('audit_id', 'N/A')}",
            f"**Date:** {data.get('timestamp', 'N/A')}",
            f"**Duration:** {data.get('duration', 0):.1f} seconds",
            f"**Target:** Healthcare Clinical Documentation System",
            "",
            "## Executive Summary",
            critique.get("summary", "No summary available."),
            "",
            "## Key Metrics",
            f"- Overall Risk: **{risk.get('overall_risk', 'UNKNOWN')}**",
            f"- Threats Detected: **{risk.get('threats_detected', 0)}**",
            f"- HIPAA Compliance Score: **{compliance.get('score', 0):.1f}%**",
            f"- Fixes Applied: **{remediation.get('fixes_applied', 0)}**",
            f"- Remediation Success Rate: **{remediation.get('success_rate', 0) * 100:.1f}%**",
            "",
            "## Red Team Attack Details",
            "| Attack Type | Attempts | Successful | Success Rate |",
            "|---|---:|---:|---:|",
        ]

        if red_results:
            for atk, d in red_results.items():
                lines.append(f"| {atk.replace('_', ' ').title()} | {d.get('total_attempts', 0)} | {d.get('successful_attacks', 0)} | {d.get('success_rate', 0) * 100:.1f}% |")
        else:
            lines.append("| No attacks recorded | 0 | 0 | 0.0% |")

        lines += [
            "",
            "## Blue Team Detection",
            f"- Total threats: {blue.get('threat_count', 0)}",
            f"- Risk level: {blue.get('risk_level', 'LOW')}",
            "",
            "### Detected Anomalies",
        ]
        anomalies = blue.get("anomalies", [])
        lines += [f"- {item}" for item in anomalies] or ["- No anomalies recorded."]

        lines += [
            "",
            "## Green Team Auto-Remediation",
            f"Fixes applied: {green.get('fixes_applied', 0)}",
        ]
        for fix in green.get("fixes", []):
            status = "Applied" if fix.get("fix_applied") else "Not applied"
            lines.append(f"- {status}: {fix.get('fix_description', 'N/A')}")
        if not green.get("fixes"):
            lines.append("- No remediation actions recorded.")

        lines += [
            "",
            "## HIPAA Compliance Validation",
            f"- Standard: {comp_details.get('standard', 'HIPAA')}",
            f"- Compliance Score: {comp_details.get('score', 0) * 100:.1f}%",
            f"- Checks Passed: {comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}",
            f"- Status: {'COMPLIANT' if comp_details.get('compliant', False) else 'NON-COMPLIANT'}",
            "",
            "## Recommendations",
        ]
        for i, rec in enumerate(data.get("recommendations", []), 1):
            lines.append(f"{i}. {rec}")

        narrative = "\n".join(lines)

        # Save
        draft_path = os.path.join(self.reports_dir, "report_draft.md")
        with open(draft_path, "w", encoding="utf-8") as f:
            f.write(narrative)
        return narrative

    # ════════════════════════════════════════════════════════════════════════
    # PHASE 5: RENDER AGENT — Multi-format export
    # ════════════════════════════════════════════════════════════════════════

    def _run_render_agent(self, narrative: str, visuals: Dict, critique: Dict) -> Dict[str, str]:
        outputs = {}
        formats_to_generate = {
            "pdf":  self._export_pdf,
            "docx": self._export_docx,
            "xlsx": self._export_xlsx,
            "pptx": self._export_pptx,
            "html": self._export_html,
            "json": self._export_json,
        }
        for fmt, fn in formats_to_generate.items():
            try:
                path = fn(narrative, visuals, critique)
                if path:
                    outputs[fmt] = path
            except Exception as e:
                print(f"   {fmt.upper()} export failed: {e}")
        return outputs

    # ── Helper: image to data URI ───────────────────────────────────────────
    def _image_to_data_uri(self, image_path: Optional[str]) -> str:
        if not image_path or not os.path.exists(image_path):
            return ""
        try:
            with open(image_path, "rb") as f:
                encoded = base64.b64encode(f.read()).decode("ascii")
            return f"data:image/png;base64,{encoded}"
        except Exception:
            return ""

    @staticmethod
    def _esc(value: Any) -> str:
        return str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    # ════════════════════════════════════════════════════════════════════════
    # PDF EXPORT (reportlab — professional enterprise-grade)
    # ════════════════════════════════════════════════════════════════════════

    def _export_pdf(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        try:
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.platypus import (
                Image, KeepTogether, Paragraph, SimpleDocTemplate,
                Spacer, Table, TableStyle, PageBreak,
            )
        except ImportError:
            print("   reportlab not installed — skipping PDF")
            return None

        data = self.report_data
        context = data.get("report_context", {})
        risk = data.get("risk_assessment", {})
        remediation = data.get("remediation", {})
        compliance = data.get("compliance", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})
        phi = data.get("phi_anonymization", {})
        baseline = data.get("baseline", {})

        institution = context.get("institution_name", "Demo Healthcare Institution")
        report_date = context.get("report_generated_date", datetime.now().date())
        target = context.get("target_system", "Healthcare Clinical Documentation System")

        navy_c = colors.HexColor(self.NAVY)
        teal_c = colors.HexColor(self.TEAL)
        red_c = colors.HexColor(self.RED)
        gold_c = colors.HexColor(self.GOLD)
        green_c = colors.HexColor(self.GREEN)
        soft_c = colors.HexColor(self.SOFT)
        line_c = colors.HexColor(self.LINE)
        ink_c = colors.HexColor(self.INK)
        light_gold_c = colors.HexColor(self.LIGHT_GOLD)
        light_teal_c = colors.HexColor(self.LIGHT_TEAL)

        path = os.path.join(self.reports_dir, "audit_report.pdf")
        doc = SimpleDocTemplate(
            str(path), pagesize=letter,
            title="Med-Sec Audit Agent Report",
            rightMargin=0.55 * inch, leftMargin=0.55 * inch,
            topMargin=0.46 * inch, bottomMargin=0.66 * inch,
        )
        story: List[Any] = []

        # ── Styles ──────────────────────────────────────────────────────────
        styles = getSampleStyleSheet()
        
        # Cover styles
        styles.add(ParagraphStyle(name="CoverTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=26, leading=30, textColor=colors.white, alignment=TA_CENTER, spaceAfter=6))
        styles.add(ParagraphStyle(name="CoverSub", parent=styles["BodyText"], fontSize=13, leading=16, textColor=light_teal_c, alignment=TA_CENTER, spaceAfter=8))
        styles.add(ParagraphStyle(name="CoverTagline", parent=styles["BodyText"], fontName="Helvetica-Bold", fontSize=12, leading=14, textColor=light_gold_c, alignment=TA_CENTER, spaceAfter=10))
        styles.add(ParagraphStyle(name="CoverDesc", parent=styles["BodyText"], fontSize=9, leading=11, textColor=light_teal_c, alignment=TA_CENTER, spaceAfter=6))
        
        # Body styles
        styles.add(ParagraphStyle(name="Section", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=15, leading=18, textColor=navy_c, spaceBefore=7, spaceAfter=6))
        styles.add(ParagraphStyle(name="SubSection", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=11.5, leading=14, textColor=teal_c, spaceBefore=5, spaceAfter=4))
        styles.add(ParagraphStyle(name="BodySmall", parent=styles["BodyText"], fontSize=9.2, leading=12.2, textColor=ink_c))
        styles.add(ParagraphStyle(name="Cell", parent=styles["BodyText"], fontSize=8.4, leading=10.5, textColor=ink_c))
        styles.add(ParagraphStyle(name="WhiteCell", parent=styles["BodyText"], fontSize=8.5, leading=10.5, textColor=colors.white))
        styles.add(ParagraphStyle(name="MetricLabel", parent=styles["BodyText"], fontSize=8, leading=10, textColor=colors.white, alignment=TA_CENTER))
        styles.add(ParagraphStyle(name="MetricValue", parent=styles["BodyText"], fontName="Helvetica-Bold", fontSize=17, leading=20, textColor=colors.white, alignment=TA_CENTER))

        from html import escape as _esc_html
        def para(value, style="Cell"):
            return Paragraph(_esc_html(str(value)), styles[style])

        def styled_table(tbl_data, widths, header=True):
            t = Table(tbl_data, colWidths=widths, hAlign="LEFT", repeatRows=1 if header else 0)
            cmds = [
                ("BOX", (0, 0), (-1, -1), 0.6, line_c),
                ("INNERGRID", (0, 0), (-1, -1), 0.35, line_c),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 7),
                ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ("ROWBACKGROUNDS", (0, 1 if header else 0), (-1, -1), [colors.white, soft_c]),
            ]
            if header:
                cmds += [("BACKGROUND", (0, 0), (-1, 0), navy_c), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold")]
            t.setStyle(TableStyle(cmds))
            return t

        def image_flowable(image_path: str, max_w=6.6 * inch, max_h=3.15 * inch):
            if not image_path or not os.path.exists(image_path):
                return None
            try:
                img = Image(image_path)
                # Use _restrictSize so reportlab knows the exact rendered dimensions
                # during the wrapping/layout phase — this prevents blank-page gaps.
                img._restrictSize(max_w, max_h)
                img.hAlign = 'CENTER'
                return img
            except Exception:
                return None

        def footer(canvas, d):
            canvas.saveState()
            canvas.setStrokeColor(teal_c)
            canvas.setLineWidth(1.1)
            canvas.line(d.leftMargin, 0.55 * inch, letter[0] - d.rightMargin, 0.55 * inch)
            canvas.setFont("Helvetica", 8)
            canvas.setFillColor(colors.HexColor("#667085"))
            canvas.drawString(d.leftMargin, 0.35 * inch, "Med-Sec Audit Agent | Confidential security assessment")
            canvas.drawRightString(letter[0] - d.rightMargin, 0.35 * inch, f"Page {d.page}")
            canvas.restoreState()

        # ── Cover Page ──────────────────────────────────────────────────────
        cover_title = Table(
            [
                [Paragraph(_esc_html(str(institution)), styles["CoverTitle"])],
                [Paragraph("Med-Sec Audit Agent Security & Compliance Audit Report", styles["CoverSub"])],
                [Paragraph("(Med-Sec Audit Agent: AI-Powered Healthcare Clinical Data Security & Compliance Auditor)", styles["CoverDesc"])],
                [Spacer(1, 4)],
                [Paragraph("DATA SECURITY IS NOT A FEATURE — IT IS A NECESSITY.<br/>AND NOW, IT IS AUTOMATED.", styles["CoverTagline"])],
                [Spacer(1, 4)],
                [Paragraph("TARGET SYSTEM: Healthcare Clinical Data Documentation System (EHR)", styles["CoverDesc"])],
            ],
            colWidths=[6.8 * inch],
            rowHeights=[0.55 * inch, 0.28 * inch, 0.20 * inch, 0.06 * inch, 0.35 * inch, 0.06 * inch, 0.22 * inch],
        )
        cover_title.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), navy_c),
            ("BOX", (0, 0), (-1, -1), 0, navy_c),
            ("LEFTPADDING", (0, 0), (-1, -1), 24),
            ("RIGHTPADDING", (0, 0), (-1, -1), 24),
            ("TOPPADDING", (0, 0), (-1, -1), 16),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 12),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(cover_title)
        story.append(Spacer(1, 8))

        # ── Cover Meta Card ─────────────────────────────────────────────────
        cl = ParagraphStyle(name="CL", parent=styles["BodyText"], fontSize=7.8, leading=9.5, textColor=colors.HexColor("#5A6B7D"), spaceAfter=1)
        cv = ParagraphStyle(name="CV", parent=styles["BodyText"], fontName="Helvetica-Bold", fontSize=10.2, leading=12.5, textColor=navy_c)
        cover_meta = [
            [Paragraph("REPORT DATE", cl), Paragraph("TARGET SYSTEM", cl)],
            [Paragraph(_esc_html(str(report_date)), cv), Paragraph(_esc_html("Healthcare Clinical Data Documentation System (EHR)"), cv)],
            [Paragraph("AUDIT ID", cl), Paragraph("GENERATED BY", cl)],
            [Paragraph(_esc_html(str(data.get("audit_id", "N/A"))), cv), Paragraph("Med-Sec Audit Agent (<b>PHI Guardians</b>)", cv)],
        ]
        meta_card = Table(cover_meta, colWidths=[3.25 * inch, 3.25 * inch], rowHeights=[0.14 * inch, 0.28 * inch, 0.14 * inch, 0.28 * inch])
        meta_card.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), soft_c),
            ("BOX", (0, 0), (-1, -1), 0.8, line_c),
            ("INNERGRID", (0, 0), (-1, -1), 0.3, line_c),
            ("LEFTPADDING", (0, 0), (-1, -1), 12),
            ("RIGHTPADDING", (0, 0), (-1, -1), 12),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LINEBEFORE", (1, 0), (1, -1), 1.2, teal_c),
        ]))
        story.append(meta_card)
        story.append(Spacer(1, 8))

        # ── Scope Strip ─────────────────────────────────────────────────────
        scope_tile = ParagraphStyle(name="ST", parent=styles["BodyText"], alignment=TA_CENTER, fontName="Helvetica-Bold", fontSize=13.2, leading=15, textColor=colors.white)
        scope_strip = Table(
            [[Paragraph("Multi-Agent Audit", scope_tile), Paragraph("HIPAA Validation", scope_tile), Paragraph("PHI Protection", scope_tile)]],
            colWidths=[2.16 * inch, 2.16 * inch, 2.16 * inch],
            rowHeights=[0.40 * inch],
        )
        scope_strip.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, 0), teal_c),
            ("BACKGROUND", (1, 0), (1, 0), red_c),
            ("BACKGROUND", (2, 0), (2, 0), gold_c),
            ("INNERGRID", (0, 0), (-1, -1), 4, colors.white),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(scope_strip)
        story.append(Spacer(1, 10))

        # ── Executive Summary ───────────────────────────────────────────────
        story.append(Paragraph("Executive Summary", styles["Section"]))
        story.append(Paragraph(
            f"This assessment processed the Healthcare Clinical Documentation System against adversarial testing, "
            f"threat detection, automated remediation, PHI masking, and HIPAA-oriented compliance validation. "
            f"The audit identified a <b>{risk.get('overall_risk', 'UNKNOWN')}</b> risk level with "
            f"<b>{risk.get('threats_detected', 0)}</b> threats detected. HIPAA compliance score: "
            f"<b>{compliance.get('score', 0):.1f}%</b>.",
            styles["BodySmall"],
        ))
        story.append(Spacer(1, 4))

        # ── Metric Cards ────────────────────────────────────────────────────
        metric_data = [[
            Paragraph("Overall Risk", styles["MetricLabel"]),
            Paragraph("Threats", styles["MetricLabel"]),
            Paragraph("HIPAA Score", styles["MetricLabel"]),
            Paragraph("Fixes Applied", styles["MetricLabel"]),
        ], [
            Paragraph(str(risk.get("overall_risk", "UNKNOWN")), styles["MetricValue"]),
            Paragraph(str(risk.get("threats_detected", 0)), styles["MetricValue"]),
            Paragraph(f"{compliance.get('score', 0):.1f}%", styles["MetricValue"]),
            Paragraph(str(remediation.get("fixes_applied", 0)), styles["MetricValue"]),
        ]]
        metric_table = Table(metric_data, colWidths=[1.65 * inch] * 4, rowHeights=[0.33 * inch, 0.48 * inch])
        metric_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), red_c),
            ("BACKGROUND", (1, 0), (1, -1), colors.HexColor("#8C2F39")),
            ("BACKGROUND", (2, 0), (2, -1), green_c),
            ("BACKGROUND", (3, 0), (3, -1), gold_c),
            ("BOX", (0, 0), (-1, -1), 0, colors.white),
            ("INNERGRID", (0, 0), (-1, -1), 3, colors.white),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story += [metric_table, Spacer(1, 4)]

        # ── Processing Scope ────────────────────────────────────────────────
        story.append(Paragraph("Processing Scope", styles["Section"]))
        scope = [
            [para("Data / Control Area", "WhiteCell"), para("Evidence Processed", "WhiteCell")],
            [para("Baseline status"), para(baseline.get("status", "N/A"))],
            [para("Attack categories tested"), para(len(red_results))],
            [para("PHI elements identified/masked"), para(phi.get("original_phi_count", 0))],
            [para("Compliance checks passed"), para(f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}")],
        ]
        story += [styled_table(scope, [2.15 * inch, 4.45 * inch]), Spacer(1, 4)]

        # ── Security Team Assessment ─────────────────────────────────────────
        story.append(Paragraph("Security Team Assessment", styles["Section"]))
        team_rows = [
            [para("Team", "WhiteCell"), para("Mission", "WhiteCell"), para("Outcome", "WhiteCell")],
            [para("Red Team"), para("Adversarial prompt injection, PHI extraction, jailbreak, role override, data exfiltration."), para(f"{len(red_results)} categories tested; overall risk {risk.get('overall_risk', 'UNKNOWN')}.")],
            [para("Blue Team"), para("Runtime detection of suspicious instructions, anomalous inputs, policy violations."), para(f"{blue.get('threat_count', 0)} threats detected; risk level {blue.get('risk_level', 'LOW')}.")],
            [para("Green Team"), para("Automated remediation actions and control hardening."), para(f"{green.get('fixes_applied', 0)} remediation actions documented.")],
        ]
        story += [styled_table(team_rows, [1.15 * inch, 3.25 * inch, 2.2 * inch]), Spacer(1, 4)]

        # ── Visual Evidence — starts on a fresh page ────────────────────────
        story.append(PageBreak())
        story.append(Paragraph("Visual Evidence", styles["Section"]))
        story.append(Spacer(1, 2))

        # First four small charts flow together on this page
        small_charts = [
            ("Risk Level Chart", "risk_chart", 5.5 * inch, 1.6 * inch),
            ("Attack Success Rate Chart", "attack_chart", 5.5 * inch, 1.6 * inch),
            ("Compliance Dashboard", "compliance_dashboard", 5.5 * inch, 1.6 * inch),
            ("Team Summary Chart", "team_summary", 5.5 * inch, 1.6 * inch),
        ]
        for title, key, width, height in small_charts:
            img = image_flowable(visuals.get(key, ""), max_w=width, max_h=height)
            if img:
                story.append(Paragraph(title, styles["SubSection"]))
                story.append(img)
                story.append(Spacer(1, 2))

        # ── Threat Heatmap — pushed to its own fresh page ────────────────────
        story.append(PageBreak())
        heatmap_img = image_flowable(visuals.get("threat_heatmap", ""), max_w=6.5 * inch, max_h=2.8 * inch)
        if heatmap_img:
            story.append(Paragraph("Threat Heatmap", styles["SubSection"]))
            story.append(heatmap_img)
            story.append(Spacer(1, 2))

        # ── Red Team Detail — follows immediately on the same page ───────────
        story.append(Paragraph("Red Team Attack Detail", styles["Section"]))
        red_rows = [[para("Attack Type", "WhiteCell"), para("Attempts", "WhiteCell"), para("Successful", "WhiteCell"), para("Success Rate", "WhiteCell")]]
        for atk, d in red_results.items():
            red_rows.append([para(atk.replace("_", " ").title()), para(d.get("total_attempts", 0)), para(d.get("successful_attacks", 0)), para(f"{d.get('success_rate', 0) * 100:.1f}%")])
        if len(red_rows) == 1:
            red_rows.append([para("No attacks recorded"), para(0), para(0), para("0.0%")])
        story += [styled_table(red_rows, [2.65 * inch, 1.2 * inch, 1.2 * inch, 1.35 * inch]), Spacer(1, 4)]

        # ── Compliance & Remediation ─────────────────────────────────────────
        story.append(Paragraph("Compliance And Remediation", styles["Section"]))
        comp_rows = [
            [para("Control", "WhiteCell"), para("Status", "WhiteCell")],
            [para("Standard"), para(comp_details.get("standard", "HIPAA"))],
            [para("Compliance score"), para(f"{comp_details.get('score', 0) * 100:.1f}%")],
            [para("Checks passed"), para(f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}")],
            [para("Compliance status"), para("COMPLIANT" if comp_details.get("compliant", False) else "NON-COMPLIANT")],
        ]
        story += [styled_table(comp_rows, [2.15 * inch, 4.45 * inch]), Spacer(1, 3)]

        issues = comp_details.get("issues", []) or ["No compliance issues recorded."]
        story.append(Paragraph("Issues Found", styles["SubSection"]))
        for issue in issues:
            story.append(Paragraph(f"- {_esc_html(str(issue))}", styles["BodySmall"]))
        story.append(Spacer(1, 4))

        # ── Recommendations ─────────────────────────────────────────────────
        rec_block = [Paragraph("Priority Recommendations", styles["Section"])]
        for i, rec in enumerate(data.get("recommendations", []) or ["No recommendations recorded."], 1):
            rec_block.append(Paragraph(f"{i}. {_esc_html(str(rec))}", styles["BodySmall"]))
            rec_block.append(Spacer(1, 2))
        story.append(KeepTogether(rec_block))

        doc.build(story, onFirstPage=footer, onLaterPages=footer)
        print(f"   PDF report: {path}")
        return path

    # ════════════════════════════════════════════════════════════════════════
    # DOCX EXPORT (python-docx — professional Word document)
    # ════════════════════════════════════════════════════════════════════════

    def _export_docx(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        if not DOCX_AVAILABLE:
            print("   python-docx not installed — skipping DOCX")
            return None

        from docx import Document as DocxDocument
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        data = self.report_data
        context = data.get("report_context", {})
        risk = data.get("risk_assessment", {})
        remediation = data.get("remediation", {})
        compliance = data.get("compliance", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})
        phi = data.get("phi_anonymization", {})
        baseline = data.get("baseline", {})

        institution = context.get("institution_name", "Demo Healthcare Institution")
        report_date = context.get("report_generated_date", datetime.now().date())
        target = context.get("target_system", "Healthcare Clinical Documentation System")

        doc = DocxDocument()
        section = doc.sections[0]
        section.top_margin = Inches(0.65)
        section.bottom_margin = Inches(0.65)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)

        # ── Cover ────────────────────────────────────────────────────────────
        title = doc.add_heading(str(institution), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].font.color.rgb = RGBColor(0x12, 0x35, 0x5B)

        subtitle = doc.add_paragraph("Med-Sec Audit Agent Security & Compliance Audit Report")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].bold = True
        subtitle.runs[0].font.size = Pt(15)
        subtitle.runs[0].font.color.rgb = RGBColor(0x0B, 0x7A, 0x75)

        subtitle2 = doc.add_paragraph("(Med-Sec Audit Agent: AI-Powered Healthcare Clinical Data Security & Compliance Auditor)")
        subtitle2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle2.runs[0].italic = True
        subtitle2.runs[0].font.size = Pt(10)
        subtitle2.runs[0].font.color.rgb = RGBColor(0x66, 0x70, 0x85)

        tagline = doc.add_paragraph("DATA SECURITY IS NOT A FEATURE — IT IS A NECESSITY. AND NOW, IT IS AUTOMATED.")
        tagline.alignment = WD_ALIGN_PARAGRAPH.CENTER
        tagline.runs[0].bold = True
        tagline.runs[0].font.size = Pt(12)
        tagline.runs[0].font.color.rgb = RGBColor(0xB7, 0x79, 0x1F)

        target_line = doc.add_paragraph("TARGET SYSTEM: Healthcare Clinical Data Documentation System (EHR)")
        target_line.alignment = WD_ALIGN_PARAGRAPH.CENTER
        target_line.runs[0].font.size = Pt(10)
        target_line.runs[0].font.color.rgb = RGBColor(0x66, 0x70, 0x85)

        doc.add_paragraph()

        # ── Meta Table ───────────────────────────────────────────────────────
        meta_table = doc.add_table(rows=5, cols=2)
        meta_table.style = "Table Grid"
        meta_items = [
            ("Report Date", str(report_date)),
            ("Target System", "Healthcare Clinical Data Documentation System (EHR)"),
            ("Audit ID", str(data.get("audit_id", "N/A"))),
            ("Audit Timestamp", str(data.get("timestamp", "N/A"))),
            ("Generated By", context.get("prepared_by", "Med-Sec Audit Agent (PHI Guardians)")),
        ]
        for i, (label, value) in enumerate(meta_items):
            meta_table.cell(i, 0).text = label
            meta_table.cell(i, 1).text = str(value)
            for run in meta_table.cell(i, 0).paragraphs[0].runs:
                run.bold = True
                run.font.color.rgb = RGBColor(0x12, 0x35, 0x5B)

        # ── Executive Summary ────────────────────────────────────────────────
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(critique.get("summary", "No summary available."))

        # ── Risk Snapshot ────────────────────────────────────────────────────
        doc.add_heading("Enterprise Risk Snapshot", level=1)
        self._docx_callout(doc, "Overall Risk", str(risk.get("overall_risk", "UNKNOWN")), RGBColor(0xC9, 0x36, 0x36))
        self._docx_callout(doc, "Threats Detected", str(risk.get("threats_detected", 0)), RGBColor(0xC9, 0x36, 0x36))
        self._docx_callout(doc, "HIPAA Compliance Score", f"{compliance.get('score', 0):.1f}%", RGBColor(0x13, 0x79, 0x5B))
        self._docx_callout(doc, "Fixes Applied", str(remediation.get("fixes_applied", 0)), RGBColor(0xB7, 0x79, 0x1F))

        # ── Processing Scope ─────────────────────────────────────────────────
        doc.add_heading("Processing Scope", level=1)
        scope_table = doc.add_table(rows=5, cols=2)
        scope_table.style = "Table Grid"
        scope_items = [
            ("Baseline status", str(baseline.get("status", "N/A"))),
            ("Attack categories tested", str(len(red_results))),
            ("PHI elements identified/masked", str(phi.get("original_phi_count", 0))),
            ("Compliance checks passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"),
            ("Blue-team runtime", f"{blue.get('runtime_ms', 0)} ms"),
        ]
        for i, (label, value) in enumerate(scope_items):
            scope_table.cell(i, 0).text = label
            scope_table.cell(i, 1).text = value
            for run in scope_table.cell(i, 0).paragraphs[0].runs:
                run.bold = True

        # ── Security Team Assessment ───────────────────────────────────────────
        doc.add_heading("Security Team Assessment", level=1)

        doc.add_heading("Red Team — Adversarial Validation", level=2)
        doc.add_paragraph("Tested the clinical documentation workflow against prompt injection, jailbreak, PHI extraction, role override, and data exfiltration patterns.")
        red_rows = [[name.replace("_", " ").title(), str(d.get("total_attempts", 0)), str(d.get("successful_attacks", 0)), f"{d.get('success_rate', 0) * 100:.1f}%"] for name, d in red_results.items()] or [["No attacks recorded", "0", "0", "0.0%"]]
        self._docx_table(doc, ["Attack Type", "Attempts", "Successful", "Success Rate"], red_rows)

        doc.add_heading("Blue Team — Detection and Monitoring", level=2)
        doc.add_paragraph(f"The monitoring layer detected {blue.get('threat_count', 0)} threats and assigned a {blue.get('risk_level', 'LOW')} operational risk level.")
        for anomaly in blue.get("anomalies", []) or ["No anomalies recorded."]:
            doc.add_paragraph(str(anomaly), style="List Bullet")

        doc.add_heading("Green Team — Automated Remediation", level=2)
        doc.add_paragraph(f"Automated remediation completed {green.get('fixes_applied', 0)} documented control actions.")
        for fix in green.get("fixes", []) or [{"fix_description": "No remediation actions recorded.", "fix_applied": False}]:
            status = "Applied" if fix.get("fix_applied") else "Not applied"
            doc.add_paragraph(f"{status}: {fix.get('fix_description', 'N/A')}", style="List Bullet")

        # ── Visual Evidence ──────────────────────────────────────────────────
        doc.add_heading("Visual Evidence", level=1)
        for title, key in [("Risk Level Chart", "risk_chart"), ("Attack Success Rate Chart", "attack_chart"), ("Compliance Dashboard", "compliance_dashboard"), ("Team Summary", "team_summary"), ("Threat Heatmap", "threat_heatmap")]:
            img_path = visuals.get(key)
            if img_path and os.path.exists(img_path):
                doc.add_heading(title, level=2)
                # Threat heatmap gets larger size
                width = Inches(6.5) if key == "threat_heatmap" else Inches(6.3)
                doc.add_picture(img_path, width=width)

        # ── Compliance ────────────────────────────────────────────────────────
        doc.add_heading("HIPAA Compliance Validation", level=1)
        self._docx_table(doc, ["Check", "Value"], [
            ["Standard", comp_details.get("standard", "HIPAA")],
            ["Compliance Score", f"{comp_details.get('score', 0) * 100:.1f}%"],
            ["Checks Passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"],
            ["Status", "COMPLIANT" if comp_details.get("compliant", False) else "NON-COMPLIANT"],
        ])
        doc.add_heading("Issues Found", level=2)
        for issue in comp_details.get("issues", []) or ["No compliance issues recorded."]:
            doc.add_paragraph(str(issue), style="List Bullet")

        # ── Recommendations ──────────────────────────────────────────────────
        doc.add_heading("Priority Recommendations", level=1)
        for rec in data.get("recommendations", []) or ["No recommendations recorded."]:
            doc.add_paragraph(str(rec), style="List Number")

        # ── Conclusion ───────────────────────────────────────────────────────
        doc.add_heading("Management Conclusion", level=1)
        doc.add_paragraph("This report provides an executive-ready view of audit scope, adversarial testing, detection outcomes, remediation activity, compliance posture, and visual evidence in one self-contained file.")

        # ── Footer ────────────────────────────────────────────────────────────
        footer_para = doc.sections[0].footer.paragraphs[0]
        footer_para.text = "Med-Sec Audit Agent | Confidential security assessment | PHI Guardians"
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in footer_para.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x66, 0x70, 0x85)

        path = os.path.join(self.reports_dir, "audit_report.docx")
        doc.save(path)
        print(f"   DOCX report: {path}")
        return path

    def _docx_callout(self, doc, title: str, value: str, color: RGBColor):
        p = doc.add_paragraph()
        label = p.add_run(f"{title}: ")
        label.bold = True
        label.font.color.rgb = color
        val = p.add_run(str(value))
        val.bold = True

    def _docx_table(self, doc, headers: List[str], rows: List[List[str]]):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = "Table Grid"
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = str(h)
            for run in cell.paragraphs[0].runs:
                run.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            self._shade_cell(cell, "12355B")
        for row in rows:
            cells = table.add_row().cells
            for i, val in enumerate(row):
                cells[i].text = str(val)

    @staticmethod
    def _shade_cell(cell, hex_color: str):
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        shading = OxmlElement("w:shd")
        shading.set(qn("w:fill"), hex_color)
        shading.set(qn("w:val"), "clear")
        cell._tc.get_or_add_tcPr().append(shading)

    # ════════════════════════════════════════════════════════════════════════
    # EXCEL EXPORT (openpyxl — multi-sheet workbook)
    # ════════════════════════════════════════════════════════════════════════

    def _export_xlsx(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError:
            print("   openpyxl not installed — skipping XLSX")
            return None

        data = self.report_data
        context = data.get("report_context", {})
        risk = data.get("risk_assessment", {})
        remediation = data.get("remediation", {})
        compliance = data.get("compliance", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})
        phi = data.get("phi_anonymization", {})
        baseline = data.get("baseline", {})

        wb = Workbook()

        # ── Styles ───────────────────────────────────────────────────────────
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="12355B", end_color="12355B", fill_type="solid")
        title_font = Font(name="Calibri", size=16, bold=True, color="12355B")
        label_font = Font(name="Calibri", size=10, bold=True, color="172033")
        value_font = Font(name="Calibri", size=10, color="172033")
        thin_border = Border(
            left=Side(style="thin", color="D8DEE9"),
            right=Side(style="thin", color="D8DEE9"),
            top=Side(style="thin", color="D8DEE9"),
            bottom=Side(style="thin", color="D8DEE9"),
        )
        alt_fill = PatternFill(start_color="F5F7FA", end_color="F5F7FA", fill_type="solid")

        def style_header_row(ws, row, cols):
            for c in range(1, cols + 1):
                cell = ws.cell(row=row, column=c)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="left", vertical="center")
                cell.border = thin_border

        def style_data_rows(ws, start_row, end_row, cols):
            for r in range(start_row, end_row + 1):
                for c in range(1, cols + 1):
                    cell = ws.cell(row=r, column=c)
                    cell.font = value_font
                    cell.border = thin_border
                    if (r - start_row) % 2 == 1:
                        cell.fill = alt_fill

        # ── Sheet 1: Summary ──────────────────────────────────────────────────
        ws1 = wb.active
        ws1.title = "Summary"
        ws1.merge_cells("A1:B1")
        ws1["A1"] = "Med-Sec Audit Agent — Security Audit Report"
        ws1["A1"].font = title_font
        summary_items = [
            ("Institution", context.get("institution_name", "Demo Healthcare Institution")),
            ("Report Date", str(context.get("report_generated_date", ""))),
            ("Target System", "Healthcare Clinical Data Documentation System (EHR)"),
            ("Audit ID", data.get("audit_id", "N/A")),
            ("Timestamp", data.get("timestamp", "N/A")),
            ("Duration (seconds)", data.get("duration", 0)),
            ("Overall Risk", risk.get("overall_risk", "UNKNOWN")),
            ("Threats Detected", risk.get("threats_detected", 0)),
            ("HIPAA Score", f"{compliance.get('score', 0):.1f}%"),
            ("Fixes Applied", remediation.get("fixes_applied", 0)),
            ("Remediation Success Rate", f"{remediation.get('success_rate', 0) * 100:.1f}%"),
            ("PHI Items Found", phi.get("original_phi_count", 0)),
            ("Baseline Status", baseline.get("status", "N/A")),
            ("Compliance Standard", comp_details.get("standard", "HIPAA")),
            ("Compliance Score", f"{comp_details.get('score', 0) * 100:.1f}%"),
            ("Checks Passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"),
            ("Compliance Status", "COMPLIANT" if comp_details.get("compliant", False) else "NON-COMPLIANT"),
        ]
        ws1["A3"] = "Metric"
        ws1["B3"] = "Value"
        style_header_row(ws1, 3, 2)
        for i, (label, value) in enumerate(summary_items, start=4):
            ws1.cell(row=i, column=1, value=label).font = label_font
            ws1.cell(row=i, column=2, value=value).font = value_font
        style_data_rows(ws1, 4, 3 + len(summary_items), 2)
        ws1.column_dimensions["A"].width = 30
        ws1.column_dimensions["B"].width = 50

        # ── Sheet 2: Red Team ─────────────────────────────────────────────────
        ws2 = wb.create_sheet("Red Team")
        ws2.merge_cells("A1:D1")
        ws2["A1"] = "Red Team — Adversarial Attack Results"
        ws2["A1"].font = title_font
        headers2 = ["Attack Type", "Attempts", "Successful", "Success Rate (%)"]
        for i, h in enumerate(headers2, 1):
            ws2.cell(row=3, column=i, value=h)
        style_header_row(ws2, 3, 4)
        row_idx = 4
        for atk, d in red_results.items():
            ws2.cell(row=row_idx, column=1, value=atk.replace("_", " ").title())
            ws2.cell(row=row_idx, column=2, value=d.get("total_attempts", 0))
            ws2.cell(row=row_idx, column=3, value=d.get("successful_attacks", 0))
            ws2.cell(row=row_idx, column=4, value=round(d.get("success_rate", 0) * 100, 1))
            row_idx += 1
        if row_idx == 4:
            ws2.cell(row=4, column=1, value="No attacks recorded")
        style_data_rows(ws2, 4, max(row_idx - 1, 4), 4)
        for col in ["A", "B", "C", "D"]:
            ws2.column_dimensions[col].width = 22

        # ── Sheet 3: Blue Team ────────────────────────────────────────────────
        ws3 = wb.create_sheet("Blue Team")
        ws3.merge_cells("A1:C1")
        ws3["A1"] = "Blue Team — Threat Detection"
        ws3["A1"].font = title_font
        ws3["A3"] = "Metric"; ws3["B3"] = "Value"
        style_header_row(ws3, 3, 2)
        blue_items = [
            ("Threat Count", blue.get("threat_count", 0)),
            ("Risk Level", blue.get("risk_level", "LOW")),
            ("Runtime (ms)", blue.get("runtime_ms", 0)),
        ]
        for i, (label, value) in enumerate(blue_items, start=4):
            ws3.cell(row=i, column=1, value=label).font = label_font
            ws3.cell(row=i, column=2, value=value).font = value_font
        style_data_rows(ws3, 4, 3 + len(blue_items), 2)
        # Anomalies
        anomaly_start = 4 + len(blue_items) + 2
        ws3.cell(row=anomaly_start, column=1, value="Detected Anomalies").font = title_font
        ws3.cell(row=anomaly_start + 1, column=1, value="#").font = header_font
        ws3.cell(row=anomaly_start + 1, column=1).fill = header_fill
        ws3.cell(row=anomaly_start + 1, column=2, value="Anomaly").font = header_font
        ws3.cell(row=anomaly_start + 1, column=2).fill = header_fill
        anomalies = blue.get("anomalies", []) or ["No anomalies recorded."]
        for i, anom in enumerate(anomalies, start=anomaly_start + 2):
            ws3.cell(row=i, column=1, value=i - anomaly_start - 1).font = value_font
            ws3.cell(row=i, column=2, value=str(anom)).font = value_font
        ws3.column_dimensions["A"].width = 20
        ws3.column_dimensions["B"].width = 60

        # ── Sheet 4: Green Team ───────────────────────────────────────────────
        ws4 = wb.create_sheet("Green Team")
        ws4.merge_cells("A1:C1")
        ws4["A1"] = "Green Team — Auto-Remediation"
        ws4["A1"].font = title_font
        headers4 = ["#", "Status", "Fix Description"]
        for i, h in enumerate(headers4, 1):
            ws4.cell(row=3, column=i, value=h)
        style_header_row(ws4, 3, 3)
        fixes = green.get("fixes", [])
        if fixes:
            for i, fix in enumerate(fixes, start=4):
                status = "Applied" if fix.get("fix_applied") else "Not applied"
                ws4.cell(row=i, column=1, value=i - 3)
                ws4.cell(row=i, column=2, value=status)
                ws4.cell(row=i, column=3, value=fix.get("fix_description", "N/A"))
            style_data_rows(ws4, 4, 3 + len(fixes), 3)
        else:
            ws4.cell(row=4, column=1, value="No remediation actions recorded.")
        ws4.column_dimensions["A"].width = 6
        ws4.column_dimensions["B"].width = 15
        ws4.column_dimensions["C"].width = 60

        # ── Sheet 5: Compliance ───────────────────────────────────────────────
        ws5 = wb.create_sheet("Compliance")
        ws5.merge_cells("A1:B1")
        ws5["A1"] = "HIPAA Compliance Validation"
        ws5["A1"].font = title_font
        comp_items = [
            ("Standard", comp_details.get("standard", "HIPAA")),
            ("Compliance Score", f"{comp_details.get('score', 0) * 100:.1f}%"),
            ("Checks Passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"),
            ("Status", "COMPLIANT" if comp_details.get("compliant", False) else "NON-COMPLIANT"),
        ]
        ws5["A3"] = "Check"; ws5["B3"] = "Value"
        style_header_row(ws5, 3, 2)
        for i, (label, value) in enumerate(comp_items, start=4):
            ws5.cell(row=i, column=1, value=label).font = label_font
            ws5.cell(row=i, column=2, value=value).font = value_font
        style_data_rows(ws5, 4, 3 + len(comp_items), 2)
        # Issues
        issue_start = 4 + len(comp_items) + 2
        ws5.cell(row=issue_start, column=1, value="Issues Found").font = title_font
        issues = comp_details.get("issues", []) or ["No compliance issues recorded."]
        for i, issue in enumerate(issues, start=issue_start + 1):
            ws5.cell(row=i, column=1, value=f"{i - issue_start}.").font = value_font
            ws5.cell(row=i, column=2, value=str(issue)).font = value_font
        ws5.column_dimensions["A"].width = 22
        ws5.column_dimensions["B"].width = 60

        # ── Sheet 6: Recommendations ─────────────────────────────────────────
        ws6 = wb.create_sheet("Recommendations")
        ws6.merge_cells("A1:B1")
        ws6["A1"] = "Priority Recommendations"
        ws6["A1"].font = title_font
        ws6["A3"] = "#"; ws6["B3"] = "Recommendation"
        style_header_row(ws6, 3, 2)
        recs = data.get("recommendations", []) or ["No recommendations recorded."]
        for i, rec in enumerate(recs, start=4):
            ws6.cell(row=i, column=1, value=i - 3).font = value_font
            ws6.cell(row=i, column=2, value=str(rec)).font = value_font
        style_data_rows(ws6, 4, 3 + len(recs), 2)
        ws6.column_dimensions["A"].width = 6
        ws6.column_dimensions["B"].width = 70

        path = os.path.join(self.reports_dir, "audit_report.xlsx")
        wb.save(path)
        print(f"   XLSX report: {path}")
        return path

    # ════════════════════════════════════════════════════════════════════════
    # POWERPOINT EXPORT (python-pptx — professional slide deck)
    # ════════════════════════════════════════════════════════════════════════

    def _export_pptx(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            print("   python-pptx not installed — skipping PPTX")
            return None

        data = self.report_data
        context = data.get("report_context", {})
        risk = data.get("risk_assessment", {})
        remediation = data.get("remediation", {})
        compliance = data.get("compliance", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green_team_data = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})
        phi = data.get("phi_anonymization", {})
        baseline = data.get("baseline", {})

        institution = context.get("institution_name", "Demo Healthcare Institution")
        report_date = context.get("report_generated_date", datetime.now().date())
        target = context.get("target_system", "Healthcare Clinical Documentation System")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # ── Colors ────────────────────────────────────────────────────────────
        navy = RGBColor(0x12, 0x35, 0x5B)
        teal = RGBColor(0x0B, 0x7A, 0x75)
        red = RGBColor(0xC9, 0x36, 0x36)
        gold = RGBColor(0xB7, 0x79, 0x1F)
        green = RGBColor(0x13, 0x79, 0x5B)
        white = RGBColor(0xFF, 0xFF, 0xFF)
        ink = RGBColor(0x17, 0x20, 0x33)
        muted = RGBColor(0x66, 0x70, 0x85)
        soft = RGBColor(0xF5, 0xF7, 0xFA)
        light_gold = RGBColor(0xF5, 0xE6, 0xB8)
        light_teal = RGBColor(0xBD, 0xE7, 0xE3)

        def add_bg_rect(slide, color, left=0, top=0, width=prs.slide_width, height=prs.slide_height):
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=ink, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = anchor
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = align
            return txBox

        def add_table(slide, left, top, width, height, headers, rows):
            num_rows = len(rows) + 1
            num_cols = len(headers)
            table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
            table = table_shape.table
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(h)
                cell.fill.solid()
                cell.fill.fore_color.rgb = navy
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = white
            for r_idx, row in enumerate(rows, 1):
                for c_idx, val in enumerate(row):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = soft if r_idx % 2 == 0 else white
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = ink
            return table_shape

        # ── Slide 1: Cover ───────────────────────────────────────────────────
        slide1 = prs.slides.add_slide(blank_layout)
        add_bg_rect(slide1, navy)
        add_textbox(slide1, Inches(1), Inches(1.2), Inches(11.3), Inches(1.3), str(institution), 34, True, white, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_textbox(slide1, Inches(1), Inches(2.6), Inches(11.3), Inches(0.7), "Med-Sec Audit Agent Security & Compliance Audit Report", 18, False, light_teal, PP_ALIGN.CENTER)
        add_textbox(slide1, Inches(1), Inches(3.3), Inches(11.3), Inches(0.5), "(Med-Sec Audit Agent: AI-Powered Healthcare Clinical Data Security & Compliance Auditor)", 11, False, muted, PP_ALIGN.CENTER)
        add_textbox(slide1, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8), "DATA SECURITY IS NOT A FEATURE — IT IS A NECESSITY.\nAND NOW, IT IS AUTOMATED.", 14, True, light_gold, PP_ALIGN.CENTER)
        add_textbox(slide1, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4), "TARGET SYSTEM: Healthcare Clinical Data Documentation System (EHR)", 12, False, light_teal, PP_ALIGN.CENTER)
        add_textbox(slide1, Inches(1), Inches(6.2), Inches(11.3), Inches(0.4), f"Report Date: {report_date}  |  Audit ID: {data.get('audit_id', 'N/A')}", 11, False, muted, PP_ALIGN.CENTER)
        add_textbox(slide1, Inches(1), Inches(6.8), Inches(11.3), Inches(0.3), "Prepared by Med-Sec Audit Agent (PHI Guardians)", 10, False, muted, PP_ALIGN.CENTER)

        # ── Slide 2: Executive Summary ────────────────────────────────────────
        slide2 = prs.slides.add_slide(blank_layout)
        add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Executive Summary", 28, True, navy)
        # Metric cards
        metric_y = Inches(1.2)
        metrics = [
            ("Overall Risk", str(risk.get("overall_risk", "UNKNOWN")), red),
            ("Threats Detected", str(risk.get("threats_detected", 0)), red),
            ("HIPAA Score", f"{compliance.get('score', 0):.1f}%", green),
            ("Fixes Applied", str(remediation.get("fixes_applied", 0)), gold),
        ]
        for i, (label, value, color) in enumerate(metrics):
            x = Inches(0.5 + i * 3.1)
            shape = slide2.shapes.add_shape(1, x, metric_y, Inches(2.8), Inches(1.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(12)
            p1.font.color.rgb = white
            p1.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = white
            p2.alignment = PP_ALIGN.CENTER
        add_textbox(slide2, Inches(0.5), Inches(3.2), Inches(12), Inches(3), critique.get("summary", "No summary available."), 14, False, ink)

        # ── Slide 3: Processing Scope ─────────────────────────────────────────
        slide3 = prs.slides.add_slide(blank_layout)
        add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Processing Scope", 28, True, navy)
        scope_rows = [
            ["Baseline status", str(baseline.get("status", "N/A"))],
            ["Attack categories tested", str(len(red_results))],
            ["PHI elements identified/masked", str(phi.get("original_phi_count", 0))],
            ["Compliance checks passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"],
            ["Blue-team runtime", f"{blue.get('runtime_ms', 0)} ms"],
        ]
        add_table(slide3, Inches(0.5), Inches(1.2), Inches(12), Inches(3), ["Data / Control Area", "Evidence Processed"], scope_rows)

        # ── Slide 4: Red Team Attack Detail ────────────────────────────────────
        slide4 = prs.slides.add_slide(blank_layout)
        add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Red Team — Attack Detail", 28, True, navy)
        red_rows = []
        for atk, d in red_results.items():
            red_rows.append([atk.replace("_", " ").title(), str(d.get("total_attempts", 0)), str(d.get("successful_attacks", 0)), f"{d.get('success_rate', 0) * 100:.1f}%"])
        if not red_rows:
            red_rows = [["No attacks recorded", "0", "0", "0.0%"]]
        add_table(slide4, Inches(0.5), Inches(1.2), Inches(12), Inches(3), ["Attack Type", "Attempts", "Successful", "Success Rate"], red_rows)

        # ── Slide 5: Visual Evidence ──────────────────────────────────────────
        slide5 = prs.slides.add_slide(blank_layout)
        add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Visual Evidence", 28, True, navy)
        chart_images = [
            ("risk_chart", Inches(0.5), Inches(1.1), Inches(6), Inches(2.8)),
            ("attack_chart", Inches(6.8), Inches(1.1), Inches(6), Inches(2.8)),
            ("compliance_dashboard", Inches(0.5), Inches(4.1), Inches(6), Inches(2.8)),
            ("team_summary", Inches(6.8), Inches(4.1), Inches(6), Inches(2.8)),
            ("threat_heatmap", Inches(0.5), Inches(6.9), Inches(12), Inches(0.6)),
        ]
        for key, left, top, width, height in chart_images:
            img_path = visuals.get(key)
            if img_path and os.path.exists(img_path):
                try:
                    slide5.shapes.add_picture(img_path, left, top, width=width, height=height)
                except Exception:
                    pass

        # ── Slide 6: Security Team Assessment ─────────────────────────────────
        slide6 = prs.slides.add_slide(blank_layout)
        add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Security Team Assessment", 28, True, navy)
        team_data = [
            ["Red Team", "Adversarial Testing", f"{len(red_results)} categories; risk {risk.get('overall_risk', 'UNKNOWN')}"],
            ["Blue Team", "Threat Detection", f"{blue.get('threat_count', 0)} threats; {blue.get('risk_level', 'LOW')} risk"],
            ["Green Team", "Auto-Remediation", f"{green_team_data.get('fixes_applied', 0)} fixes applied"],
        ]
        add_table(slide6, Inches(0.5), Inches(1.2), Inches(12), Inches(2.5), ["Team", "Mission", "Outcome"], team_data)

        # ── Slide 7: Compliance ───────────────────────────────────────────────
        slide7 = prs.slides.add_slide(blank_layout)
        add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "HIPAA Compliance Validation", 28, True, navy)
        comp_rows = [
            ["Standard", str(comp_details.get("standard", "HIPAA"))],
            ["Compliance Score", f"{comp_details.get('score', 0) * 100:.1f}%"],
            ["Checks Passed", f"{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}"],
            ["Status", "COMPLIANT" if comp_details.get("compliant", False) else "NON-COMPLIANT"],
        ]
        add_table(slide7, Inches(0.5), Inches(1.2), Inches(12), Inches(2.5), ["Check", "Value"], comp_rows)
        issues = comp_details.get("issues", [])
        if issues:
            add_textbox(slide7, Inches(0.5), Inches(4.2), Inches(12), Inches(2.5), "Issues Found:\n" + "\n".join(f"  • {issue}" for issue in issues[:5]), 12, False, ink)

        # ── Slide 8: Recommendations ──────────────────────────────────────────
        slide8 = prs.slides.add_slide(blank_layout)
        add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), "Priority Recommendations", 28, True, navy)
        recs = data.get("recommendations", []) or ["No recommendations recorded."]
        rec_text = "\n".join(f"{i}. {rec}" for i, rec in enumerate(recs, 1))
        add_textbox(slide8, Inches(0.5), Inches(1.2), Inches(12), Inches(5), rec_text, 14, False, ink)

        # ── Slide 9: Conclusion ────────────────────────────────────────────────
        slide9 = prs.slides.add_slide(blank_layout)
        add_bg_rect(slide9, navy)
        add_textbox(slide9, Inches(1), Inches(2.5), Inches(11.3), Inches(1), "Management Conclusion", 32, True, white, PP_ALIGN.CENTER)
        add_textbox(slide9, Inches(1), Inches(3.8), Inches(11.3), Inches(2),
                    "This report provides an executive-ready view of audit scope, adversarial testing, "
                    "detection outcomes, remediation activity, compliance posture, and visual evidence "
                    "in one self-contained package.", 16, False, light_teal, PP_ALIGN.CENTER)
        add_textbox(slide9, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4), "Med-Sec Audit Agent (PHI Guardians) | Confidential", 11, False, muted, PP_ALIGN.CENTER)

        path = os.path.join(self.reports_dir, "audit_report.pptx")
        prs.save(path)
        print(f"   PPTX report: {path}")
        return path

    # ════════════════════════════════════════════════════════════════════════
    # HTML EXPORT (self-contained, enterprise-style)
    # ════════════════════════════════════════════════════════════════════════

    def _export_html(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        data = self.report_data
        context = data.get("report_context", {})
        risk = data.get("risk_assessment", {})
        remediation = data.get("remediation", {})
        compliance = data.get("compliance", {})
        red_results = data.get("red_team", {}).get("results", {})
        blue = data.get("blue_team", {})
        green = data.get("green_team", {})
        comp_details = data.get("compliance_details", {})
        phi = data.get("phi_anonymization", {})
        baseline = data.get("baseline", {})

        institution = context.get("institution_name", "Demo Healthcare Institution")
        report_date = context.get("report_generated_date", datetime.now().date())
        target = context.get("target_system", "Healthcare Clinical Documentation System")

        metric_cards = "".join([
            f"<div class='metric risk'><span>Overall Risk</span><strong>{self._esc(risk.get('overall_risk', 'UNKNOWN'))}</strong></div>",
            f"<div class='metric red'><span>Threats Detected</span><strong>{risk.get('threats_detected', 0)}</strong></div>",
            f"<div class='metric green'><span>Compliance Score</span><strong>{compliance.get('score', 0):.1f}%</strong></div>",
            f"<div class='metric gold'><span>Fixes Applied</span><strong>{remediation.get('fixes_applied', 0)}</strong></div>",
        ])

        red_rows = "".join(
            f"<tr><td>{self._esc(name.replace('_', ' ').title())}</td>"
            f"<td>{d.get('total_attempts', 0)}</td>"
            f"<td>{d.get('successful_attacks', 0)}</td>"
            f"<td>{d.get('success_rate', 0) * 100:.1f}%</td></tr>"
            for name, d in red_results.items()
        ) or "<tr><td>No attacks recorded</td><td>0</td><td>0</td><td>0.0%</td></tr>"

        anomalies = "".join(f"<li>{self._esc(item)}</li>" for item in blue.get("anomalies", [])) or "<li>No anomalies recorded.</li>"
        fixes = "".join(
            f"<li><strong>{'Applied' if fix.get('fix_applied') else 'Not applied'}:</strong> {self._esc(fix.get('fix_description', 'N/A'))}</li>"
            for fix in green.get("fixes", [])
        ) or "<li>No remediation actions recorded.</li>"
        issues = "".join(f"<li>{self._esc(item)}</li>" for item in comp_details.get("issues", [])) or "<li>No compliance issues recorded.</li>"
        recommendations = "".join(f"<li>{self._esc(item)}</li>" for item in data.get("recommendations", [])) or "<li>No recommendations recorded.</li>"

        visual_html = ""
        for title, key in [("Risk Level", "risk_chart"), ("Attack Success Rates", "attack_chart"), ("Compliance Dashboard", "compliance_dashboard"), ("Team Summary", "team_summary"), ("Threat Heatmap", "threat_heatmap")]:
            uri = self._image_to_data_uri(visuals.get(key))
            if uri:
                visual_html += f"<figure><img src='{uri}' alt='{title}'><figcaption>{title}</figcaption></figure>"

        html_content = f"""
        <!DOCTYPE html><html><head><meta charset="UTF-8"><title>Med-Sec Audit Agent Report</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 0; color: {self.INK}; background: #EEF2F6; line-height: 1.48; }}
            .page {{ max-width: 980px; margin: 32px auto; background: white; box-shadow: 0 12px 34px rgba(23,32,51,.12); }}
            .cover {{ background: {self.NAVY}; color: white; padding: 34px 42px; border-bottom: 8px solid {self.TEAL}; }}
            .cover h1 {{ margin: 0; font-size: 30px; }} .cover h2 {{ margin: 8px 0 0; color: #BDE7E3; font-weight: 500; }}
            .cover .tagline {{ color: #F5E6B8; font-weight: bold; font-size: 14px; margin: 10px 0; }}
            .meta-grid {{ display: grid; grid-template-columns: repeat(2, 1fr); gap: 10px; margin-top: 22px; }}
            .meta {{ background: rgba(255,255,255,.10); padding: 10px 12px; border-left: 4px solid #BDE7E3; }}
            .content {{ padding: 34px 42px; }} h2 {{ color: {self.NAVY}; margin-top: 28px; border-bottom: 2px solid {self.LINE}; padding-bottom: 7px; }}
            .summary {{ background: {self.SOFT}; border-left: 6px solid {self.TEAL}; padding: 16px 18px; margin: 14px 0; }}
            .metrics {{ display:grid; grid-template-columns: repeat(4, 1fr); gap: 12px; margin: 18px 0 24px; }}
            .metric {{ border-radius: 8px; padding: 14px; color: white; min-height: 78px; }} .metric span {{ display:block; font-size: 12px; opacity:.88; }} .metric strong {{ display:block; font-size: 25px; margin-top: 6px; }}
            .risk {{ background:#7A1E1E; }} .red {{ background:{self.RED}; }} .green {{ background:{self.GREEN}; }} .gold {{ background:{self.GOLD}; }}
            .team-grid {{ display:grid; grid-template-columns: repeat(3, 1fr); gap: 14px; }} .team {{ border:1px solid {self.LINE}; border-top: 5px solid {self.TEAL}; padding: 14px; border-radius:8px; background:#FBFCFE; }}
            table {{ border-collapse: collapse; width: 100%; margin: 14px 0 22px; }} th {{ background: {self.NAVY}; color: white; text-align: left; }} th, td {{ border: 1px solid {self.LINE}; padding: 9px; }}
            figure {{ margin: 18px 0 24px; page-break-inside: avoid; }} img {{ max-width: 100%; border: 1px solid {self.LINE}; }} figcaption {{ font-weight: bold; margin-top: 6px; color: #253244; }}
        </style></head><body><div class="page">
            <section class="cover">
                <h1>{self._esc(institution)}</h1>
                <h2>Med-Sec Audit Agent Security & Compliance Audit Report</h2>
                <p style="font-size:12px; color:#8899AA;">(Med-Sec Audit Agent: AI-Powered Healthcare Clinical Data Security & Compliance Auditor)</p>
                <div class="tagline">DATA SECURITY IS NOT A FEATURE — IT IS A NECESSITY.<br/>AND NOW, IT IS AUTOMATED.</div>
                <p style="font-size:11px; color:#BDE7E3;">TARGET SYSTEM: Healthcare Clinical Data Documentation System (EHR)</p>
                <div class="meta-grid">
                    <div class="meta"><strong>Report Date</strong><br>{self._esc(report_date)}</div>
                    <div class="meta"><strong>Target System</strong><br>Healthcare Clinical Data Documentation System (EHR)</div>
                    <div class="meta"><strong>Audit ID</strong><br>{self._esc(data.get('audit_id', 'N/A'))}</div>
                    <div class="meta"><strong>Generated By</strong><br>{self._esc(context.get('prepared_by', 'Med-Sec Audit Agent'))}</div>
                </div>
            </section>
            <main class="content">
                <h2>Executive Summary</h2><div class="summary">{self._esc(critique.get('summary', 'No summary available.'))}</div>
                <div class="metrics">{metric_cards}</div>
                <h2>Processing Scope</h2>
                <table><tbody>
                    <tr><td>Baseline status</td><td>{self._esc(baseline.get('status', 'N/A'))}</td></tr>
                    <tr><td>Attack categories tested</td><td>{len(red_results)}</td></tr>
                    <tr><td>PHI items identified for masking</td><td>{phi.get('original_phi_count', 0)}</td></tr>
                    <tr><td>Compliance checks passed</td><td>{comp_details.get('passed_checks', 0)}/{comp_details.get('total_checks', 0)}</td></tr>
                </tbody></table>
                <h2>Security Team Assessment</h2><div class="team-grid">
                    <div class="team"><h3>Red Team</h3><p>Validated prompt-injection, PHI-extraction, jailbreak, and data-exfiltration exposure using adversarial payloads.</p></div>
                    <div class="team"><h3>Blue Team</h3><p>Monitored attack signals and anomaly patterns. Detected {blue.get('threat_count', 0)} threats with {self._esc(blue.get('risk_level', 'LOW'))} risk.</p></div>
                    <div class="team"><h3>Green Team</h3><p>Applied remediation controls and documented fix outcomes. Fixes applied: {green.get('fixes_applied', 0)}.</p></div>
                </div>
                <h2>Visual Evidence</h2>{visual_html or '<p>No charts were generated.</p>'}
                <h2>Red Team Attack Details</h2><table><thead><tr><th>Attack Type</th><th>Attempts</th><th>Successful</th><th>Success Rate</th></tr></thead><tbody>{red_rows}</tbody></table>
                <h2>Blue Team Findings</h2><ul>{anomalies}</ul>
                <h2>Green Team Remediation</h2><ul>{fixes}</ul>
                <h2>HIPAA Compliance Validation</h2><table><tbody><tr><td>Standard</td><td>{self._esc(comp_details.get('standard', 'HIPAA'))}</td></tr><tr><td>Compliance Score</td><td>{comp_details.get('score', 0) * 100:.1f}%</td></tr><tr><td>Status</td><td>{'COMPLIANT' if comp_details.get('compliant', False) else 'NON-COMPLIANT'}</td></tr></tbody></table>
                <h3>Issues Found</h3><ul>{issues}</ul>
                <h2>Priority Recommendations</h2><ol>{recommendations}</ol>
            </main></div></body></html>
        """
        path = os.path.join(self.reports_dir, "report.html")
        with open(path, "w", encoding="utf-8") as f:
            f.write(html_content)
        print(f"   HTML report: {path}")
        return path

    # ════════════════════════════════════════════════════════════════════════
    # JSON EXPORT
    # ════════════════════════════════════════════════════════════════════════

    def _export_json(self, narrative: str, visuals: Dict, critique: Dict) -> Optional[str]:
        path = os.path.join(self.reports_dir, "audit_report.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump(self.audit_results, f, indent=2, default=str)
        print(f"   JSON report: {path}")
        return path
